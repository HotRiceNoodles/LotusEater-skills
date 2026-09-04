# -*- coding: utf-8 -*-
"""把一个目录里的 PNG 截图按页序嵌入 16:9 PPTX（每页整幅铺满）。

用法:
    python make_pptx_bitmap.py <png_dir> <out.pptx> [--pattern "p*.png"] [--missing "p00,p02-p17"]

- png_dir  : 逐页截图目录（1280x720 或任意 16:9）
- out.pptx : 输出文件
- pattern  : glob 模式，默认 p*.png
- missing  : 页名缺失时改为按自然排序全部匹配文件；也可显式给页序列表

依赖: python-pptx（venv 内 pip install python-pptx）
"""
import argparse
import os
import re
import sys


def natural_key(s):
    return [int(t) if t.isdigit() else t for t in re.split(r"(\d+)", s)]


def expand_missing(spec):
    """'p00,p02-p17' -> ['p00','p02',...,'p17']；无参数返回 None"""
    if not spec:
        return None
    names = []
    for part in spec.split(","):
        part = part.strip()
        m = re.match(r"^(.*?)(\d+)-(.*?)(\d+)$", part)
        if m and m.group(1) == m.group(3):
            w = len(m.group(2))
            names += [f"{m.group(1)}{i:0{w}d}" for i in range(int(m.group(2)), int(m.group(4)) + 1)]
        else:
            names.append(part)
    return names


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("png_dir")
    ap.add_argument("out_pptx")
    ap.add_argument("--pattern", default="p*.png")
    ap.add_argument("--missing", default=None, help="页名序列，如 p00,p02-p17；缺省时用目录内全部匹配文件自然排序")
    args = ap.parse_args()

    from pptx import Presentation
    from pptx.util import Inches

    names = expand_missing(args.missing)
    if names is None:
        files = sorted(f for f in os.listdir(args.png_dir) if f.lower().endswith(".png"))
        if not files:
            sys.exit("no png found")
        names = [os.path.splitext(f)[0] for f in files]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for name in names:
        path = os.path.join(args.png_dir, name + ".png")
        if not os.path.exists(path):
            path = os.path.join(args.png_dir, name + ".PNG")
        assert os.path.exists(path), "missing page image: " + path
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    os.makedirs(os.path.dirname(os.path.abspath(args.out_pptx)), exist_ok=True)
    prs.save(args.out_pptx)
    print("saved:", args.out_pptx, os.path.getsize(args.out_pptx) // 1024, "KB, slides:", len(names))


if __name__ == "__main__":
    main()
