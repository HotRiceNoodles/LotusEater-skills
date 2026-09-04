# -*- coding: utf-8 -*-
"""跨平台给 PPTX 注入进入动画（纯 XML 注入，不依赖 PowerPoint，Windows/macOS/Linux 通用）。

原理：直接向每页 slide XML 写出 PowerPoint 规范的 <p:timing> 树（结构与 PowerPoint
自身保存的动画 XML 一致）。python-pptx 不支持动画 API，但支持 lxml 级操作。

用法:
    python add_animations.py --src 可编辑版.pptx --dst 动画版.pptx [选项]

选项:
    --effect {fade,wipe}   默认效果（fade）
    --wide-wipe / --no-wide-wipe   宽形状(>360pt)自动改用 wipe（默认开）
    --first-duration 0.5   首个形状时长（秒）
    --duration 0.35        其余形状时长（秒）
    --stagger 0.12         级联间隔（秒），延迟封顶 2.5s
    --trigger {auto,click} auto=首形状翻页自动播+其余级联（默认）；click=全部点击触发
    --keep-pictures        背景图片也加动画（默认跳过）

动画规范备忘（与 PowerPoint COM 一致）:
    presetID: appear=1 fly=2 fade=10 wipe=22 zoom=23
    filter:   "fade" / "wipe(down)" / "fly" …
    nodeType: clickEffect=点击 / withEffect=与上一项同时 / afterEffect=上一项之后

依赖: python-pptx（自带 lxml）
"""
import argparse
import os

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

EFFECTS = {
    "fade": dict(preset=10, subtype=0, filt="fade"),
    "wipe": dict(preset=22, subtype=4, filt="wipe(down)"),
    "fly": dict(preset=2, subtype=0, filt="fly"),
    "zoom": dict(preset=23, subtype=0, filt="zoom"),
    "appear": dict(preset=1, subtype=0, filt="appear"),
}


def sub(parent, tag, **attrs):
    """建 <p:tag> 子元素，attrs 按 PowerPoint 的属性名写入"""
    el = etree.SubElement(parent, qn("p:" + tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def effect_par(parent, spid, eff, dur_ms, node_type, delay_ms, cid):
    """构造一个进入动画的 <p:par> 子树（PowerPoint 规范结构），返回下一个可用 id"""
    e = EFFECTS[eff]
    par = sub(parent, "par")
    ctn = sub(par, "cTn", id=cid, presetID=e["preset"], presetClass="entr",
              presetSubtype=e["subtype"], fill="hold", grpId=0, nodeType=node_type)
    st = sub(ctn, "stCondLst")
    sub(st, "cond", delay=delay_ms)
    child = sub(ctn, "childTnLst")

    # 1) 先把形状设为 visible（动画开始前隐藏，靠此 set 生效）
    s = sub(child, "set")
    cb = sub(s, "cBhvr")
    ctn6 = sub(cb, "cTn", id=cid + 1, dur=1, fill="hold")
    sub(sub(ctn6, "stCondLst"), "cond", delay=0)
    sub(sub(cb, "tgtEl"), "spTgt", spid=spid)
    anl = sub(cb, "attrNameLst")
    sub(anl, "attrName").text = "style.visibility"
    sub(sub(s, "to"), "strVal", val="visible")

    # 2) 进入效果本体
    ae = sub(child, "animEffect", transition="in", filter=e["filt"])
    cb2 = sub(ae, "cBhvr")
    sub(cb2, "cTn", id=cid + 2, dur=dur_ms)
    sub(sub(cb2, "tgtEl"), "spTgt", spid=spid)
    return cid + 3


def build_timing(items, first_dur, dur, stagger, trigger):
    """items = [(spid, effect)]，返回 <p:timing> 元素"""
    timing = etree.Element(qn("p:timing"))
    par0 = sub(sub(timing, "tnLst"), "par")
    ctn1 = sub(par0, "cTn", id=1, dur="indefinite", restart="never", nodeType="tmRoot")
    seq = sub(sub(ctn1, "childTnLst"), "seq", concurrent="1", nextAc="seek")
    ctn2 = sub(seq, "cTn", id=2, dur="indefinite", nodeType="mainSeq")

    par_a = sub(sub(ctn2, "childTnLst"), "par")
    ctn3 = sub(par_a, "cTn", id=3, fill="hold")
    st3 = sub(ctn3, "stCondLst")
    sub(st3, "cond", delay="indefinite")
    cond = sub(st3, "cond", evt="onBegin", delay="0")
    sub(cond, "tn", val="2")
    par_b = sub(sub(ctn3, "childTnLst"), "par")
    ctn4 = sub(par_b, "cTn", id=4, fill="hold")
    sub(sub(ctn4, "stCondLst"), "cond", delay="0")
    host = sub(ctn4, "childTnLst")

    cid = 5
    for k, (spid, eff) in enumerate(items):
        if trigger == "click":
            ctn_type, delay_ms = "clickEffect", 0
            d = first_dur if k == 0 else dur
        elif k == 0:
            ctn_type, delay_ms, d = "afterEffect", 0, first_dur   # 翻到该页自动播放
        else:
            ctn_type, delay_ms, d = "withEffect", int(min(2.5, k * stagger) * 1000), dur
        cid = effect_par(host, spid, eff, int(d * 1000), ctn_type, delay_ms, cid)

    # 幻灯片翻页事件（onPrev / onNext）
    for tag, evt in (("prevCondLst", "onPrev"), ("nextCondLst", "onNext")):
        lst = sub(seq, tag)
        c = sub(lst, "cond", evt=evt, delay="0")
        sub(sub(c, "tgtEl"), "sldTgt")

    bld = sub(timing, "bldLst")
    for spid, _ in items:
        sub(bld, "bldP", spid=spid, grpId=0)
    return timing


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", required=True)
    ap.add_argument("--dst", required=True)
    ap.add_argument("--effect", default="fade", choices=list(EFFECTS))
    ap.add_argument("--wide-wipe", dest="wide_wipe", action="store_true", default=True)
    ap.add_argument("--no-wide-wipe", dest="wide_wipe", action="store_false")
    ap.add_argument("--wipe-width", type=float, default=360.0, help="宽于此(pt)的形状改用 wipe")
    ap.add_argument("--first-duration", type=float, default=0.5)
    ap.add_argument("--duration", type=float, default=0.35)
    ap.add_argument("--stagger", type=float, default=0.12)
    ap.add_argument("--trigger", default="auto", choices=["auto", "click"])
    ap.add_argument("--keep-pictures", action="store_true")
    args = ap.parse_args()

    prs = Presentation(args.src)
    total = 0
    log = []

    for idx, slide in enumerate(prs.slides, 1):
        items = []
        for sh in slide.shapes:                      # python-pptx 按 z-order 遍历 = 生成顺序
            if not args.keep_pictures and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue                             # 背景图保持静态
            w_pt = sh.width / 12700 if sh.width else 0
            eff = "wipe" if (args.wide_wipe and w_pt > args.wipe_width) else args.effect
            items.append((sh.shape_id, eff))

        if not items:
            log.append("Slide %d: 0 effects (skipped)" % idx)
            continue

        timing = build_timing(items, args.first_duration, args.duration,
                              args.stagger, args.trigger)
        el = slide._element
        old = el.find(qn("p:timing"))
        if old is not None:
            el.remove(old)
        ext = el.find(qn("p:extLst"))
        if ext is not None:
            ext.addprevious(timing)
        else:
            el.append(timing)

        total += len(items)
        log.append("Slide %d: %d effects" % (idx, len(items)))

    os.makedirs(os.path.dirname(os.path.abspath(args.dst)), exist_ok=True)
    prs.save(args.dst)
    log.append("TOTAL=%d -> %s (%d KB)" % (total, args.dst, os.path.getsize(args.dst) // 1024))
    print("\n".join(log))


if __name__ == "__main__":
    main()
